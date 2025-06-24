import os
import tempfile
import re
from pathlib import Path

import pytest
from pptx import Presentation

from click2pptx.generate_clickable_pptx import (
    parse_html,
    create_pptx,
    find_html,
    make_output_path,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_html_text():
    path = Path(__file__).parent / "test_data" / "carte.html"
    return path.read_text(encoding="utf-8")


# ---------------------------------------------------------------------------
# parse_html
# ---------------------------------------------------------------------------


def test_parse_html_sample(sample_html_text):
    clickables, img_src = parse_html(sample_html_text)
    assert img_src == "carte.html_files/image.png"
    assert clickables == [
        (
            [140, 50, 412, 108],
            "https://news.google.com/search?q=ia&hl=fr&gl=FR&ceid=FR%3Afr",
        ),
        (
            [437, 67, 502, 91],
            "https://news.google.com/search?q=ia&hl=fr&gl=FR&ceid=FR%3Afr",
        ),
        ([100, 53, 115, 77], "https://duckduckgo.com/"),
        ([50, 81, 115, 105], "https://duckduckgo.com/"),
    ]


# ---------------------------------------------------------------------------
# create_pptx
# ---------------------------------------------------------------------------


def test_create_pptx_creates_shapes(tmp_path):
    clickables = [([10, 10, 20, 20], "http://example.com")]
    pptx_path = tmp_path / "test.pptx"
    create_pptx(clickables, None, pptx_path)

    prs = Presentation(pptx_path)
    shapes = prs.slides[0].shapes
    assert len(shapes) == 1
    assert shapes[0].click_action.hyperlink.address == "http://example.com"


# ---------------------------------------------------------------------------
# find_html
# ---------------------------------------------------------------------------


def test_find_html_existing(tmp_path):
    html_file = tmp_path / "foo.html"
    html_file.write_text("<html></html>", encoding="utf-8")
    assert find_html(str(html_file)) == str(html_file)


def test_find_html_automatic(tmp_path, monkeypatch):
    html_file = tmp_path / "bar.html"
    html_file.write_text("<html></html>", encoding="utf-8")
    monkeypatch.chdir(tmp_path)
    assert find_html(None) == html_file.name


# ---------------------------------------------------------------------------
# make_output_path
# ---------------------------------------------------------------------------


def test_make_output_path_creates_folder(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    out_path = make_output_path(None)
    assert Path(out_path).parent.exists()
    # Normalize the path to use forward slashes
    normalized_out_path = Path(out_path).as_posix()
    assert re.match(r"output/mind_map_clickable_\d{8}_\d{6}\.pptx", normalized_out_path)
