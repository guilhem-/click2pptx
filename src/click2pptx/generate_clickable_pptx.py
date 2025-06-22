#!/usr/bin/env python3
"""
generate_clickable_pptx.py
-------------------------

Convertit un export Freeplane-HTML contenant un <map> d’image‐cliquable
en un PowerPoint où chaque zone active est recréée par un rectangle
100 % transparent pointant vers le lien externe correspondant.

Usage
=====
    python generate_clickable_pptx.py [-i SOURCE.html] [-o DESTINATION.pptx]

    • -i / --input   : fichier HTML source.
                      S’il est omis, on prend le *premier* « *.html » trouvé
                      dans le répertoire courant.
    • -o / --output : chemin du fichier PPTX à produire.
                      S’il est omis, le script crée un dossier « output »
                      (s’il n’existe pas) et y écrit un fichier
                      « mind_map_clickable_YYYYMMDD_HHMMSS.pptx ».

Dépendances :  `beautifulsoup4`, `python-pptx`.
"""

import argparse
import glob
import os
import re
import sys
from datetime import datetime

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
EMU_PER_PX = 9_525  # 1 px → 9 525 EMU (pptx)

# ---------------------------------------------------------------------------#
# Outils d'E/S
# ---------------------------------------------------------------------------#

def find_html(path: str | None) -> str:
    """Retourne le chemin du fichier HTML à traiter."""
    if path:
        if not os.path.isfile(path):
            sys.exit(f"Fichier source introuvable : {path}")
        return path

    matches = glob.glob("*.html") + glob.glob("*.htm")
    if not matches:
        sys.exit("Aucun fichier *.html trouvé et aucun --input fourni.")
    return matches[0]


def make_output_path(path: str | None) -> str:
    """Construit le chemin de sortie (crée ./output si besoin)."""
    if path:
        out_dir = os.path.dirname(path)
        if out_dir and not os.path.exists(out_dir):
            os.makedirs(out_dir, exist_ok=True)
        return path

    out_dir = "output"
    os.makedirs(out_dir, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    return os.path.join(out_dir, f"mind_map_clickable_{ts}.pptx")


# ---------------------------------------------------------------------------#
# Parsing du HTML Freeplane
# ---------------------------------------------------------------------------#

def parse_html(html_text: str):
    """
    Extrait :
      • la liste des zones cliquables ([x1,y1,x2,y2], url)
      • le chemin éventuel de l'image de fond
    """
    soup = BeautifulSoup(html_text, "html.parser")

    imagemap = soup.find("map", id="fm_imagemap")
    if not imagemap:
        return [], None

    # id interne → lien externe http(s)
    id2url = {}
    for a_internal in soup.find_all("a", id=re.compile(r"FMID_\d+FM")):
        external = a_internal.find_next("a", href=re.compile(r"^https?://"))
        if external:
            id2url[a_internal["id"]] = external["href"]

    clickables = []
    for area in imagemap.find_all("area"):
        coords = list(map(int, area["coords"].split(",")))
        internal_id = area["href"].lstrip("#")
        url = id2url.get(internal_id)
        if url:
            clickables.append((coords, url))

    img_tag = soup.find("img", {"usemap": "#fm_imagemap"})
    img_src = img_tag["src"] if img_tag else None
    return clickables, img_src


# ---------------------------------------------------------------------------#
# Génération du PPTX
# ---------------------------------------------------------------------------#

def create_pptx(clickables, img_src, dest_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if img_src and os.path.isfile(img_src):
        slide.shapes.add_picture(img_src, Emu(0), Emu(0))  # image en toile de fond

    for coords, url in clickables:
        x1, y1, x2, y2 = coords
        left, top = Emu(x1 * EMU_PER_PX), Emu(y1 * EMU_PER_PX)
        width, height = Emu((x2 - x1) * EMU_PER_PX), Emu((y2 - y1) * EMU_PER_PX)

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        # REND LE RECTANGLE INVISIBLE 
        rect.fill.background()
        rect.line.fill.background()

        rect.click_action.hyperlink.address = url

    prs.save(dest_path)


# ---------------------------------------------------------------------------#
# Programme principal
# ---------------------------------------------------------------------------#

def main():
    ap = argparse.ArgumentParser(
        description="Convertit un export Freeplane-HTML en PPTX cliquable."
    )
    ap.add_argument("-i", "--input", help="Fichier HTML source")
    ap.add_argument("-o", "--output", help="Fichier PPTX de destination")
    args = ap.parse_args()

    html_path = find_html(args.input)
    output_path = make_output_path(args.output)

    with open(html_path, "r", encoding="utf-8") as f:
        html_text = f.read()

    clickables, img_src = parse_html(html_text)
    if not clickables:
        sys.exit("Aucune zone cliquable avec lien externe trouvée dans le HTML.")

    create_pptx(clickables, img_src, output_path)
    print(f"PPTX généré : {output_path}")


if __name__ == "__main__":
    main()
